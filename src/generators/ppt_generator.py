"""
PowerPoint 生成器模块

功能：
1. 基于模板生成 PPT
2. 支持多种风格（稳重、视觉、详尽）
3. 自动插入数据和图表
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
from typing import Dict, Any, List
from pathlib import Path


class StyleConfig:
    """PPT 风格配置"""
    
    CONSERVATIVE = {
        'name': '方案 A：稳重商务风',
        'primary_color': (31, 78, 120),      # 深蓝色
        'secondary_color': (79, 129, 189),   # 浅蓝色
        'background': (255, 255, 255),       # 白色
        'font_color': (0, 0, 0),             # 黑色
        'emphasis': 'data_tables'             # 强调数据表格
    }
    
    VISUAL = {
        'name': '方案 B：视觉冲击风',
        'primary_color': (255, 87, 51),      # 橙红色
        'secondary_color': (255, 195, 0),    # 金黄色
        'background': (33, 33, 33),          # 深灰色
        'font_color': (255, 255, 255),       # 白色
        'emphasis': 'large_charts'            # 强调大型图表
    }
    
    DETAILED = {
        'name': '方案 C：详尽分析风',
        'primary_color': (46, 125, 50),      # 深绿色
        'secondary_color': (129, 199, 132),  # 浅绿色
        'background': (245, 245, 245),       # 浅灰色
        'font_color': (33, 33, 33),          # 深灰色
        'emphasis': 'infographics'            # 强调信息图
    }


class PPTGenerator:
    """PowerPoint 生成器"""
    
    def __init__(self, style: str = 'conservative'):
        """
        初始化生成器
        
        Args:
            style: 风格类型 ('conservative', 'visual', 'detailed')
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # 选择风格配置
        style_map = {
            'conservative': StyleConfig.CONSERVATIVE,
            'visual': StyleConfig.VISUAL,
            'detailed': StyleConfig.DETAILED
        }
        self.style = style_map.get(style, StyleConfig.CONSERVATIVE)
        
    def add_title_slide(self, title: str, subtitle: str = None):
        """
        添加标题页
        
        Args:
            title: 主标题
            subtitle: 副标题
        """
        # 使用空白布局
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb_tuple_to_color(self.style['background'])
        
        # 添加标题
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title
        
        # 设置标题样式
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self._rgb_tuple_to_color(self.style['primary_color'])
        
        # 添加副标题
        if subtitle:
            sub_top = Inches(4.2)
            sub_box = slide.shapes.add_textbox(left, sub_top, width, Inches(0.8))
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            
            sub_p = sub_frame.paragraphs[0]
            sub_p.alignment = PP_ALIGN.CENTER
            sub_p.font.size = Pt(24)
            sub_p.font.color.rgb = self._rgb_tuple_to_color(self.style['font_color'])
        
        return slide
    
    def add_data_slide(self, title: str, df: pd.DataFrame):
        """
        添加数据表格页
        
        Args:
            title: 标题
            df: pandas DataFrame
        """
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb_tuple_to_color(self.style['background'])
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._rgb_tuple_to_color(self.style['primary_color'])
        
        # 表格
        rows, cols = min(df.shape[0] + 1, 10), min(df.shape[1], 6)  # 限制大小
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 填充表头
        for col_idx, col_name in enumerate(df.columns[:cols]):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._rgb_tuple_to_color(self.style['primary_color'])
            cell.text_frame.paragraphs[0].font.color.rgb = self._rgb_tuple_to_color(
                self.style['background']
            )
            cell.text_frame.paragraphs[0].font.bold = True
        
        # 填充数据
        for row_idx in range(min(rows - 1, df.shape[0])):
            for col_idx in range(cols):
                cell = table.cell(row_idx + 1, col_idx)
                value = df.iloc[row_idx, col_idx]
                cell.text = str(value) if pd.notna(value) else ""
                cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        return slide
    
    def add_chart_slide(self, title: str, data: Dict[str, List], chart_type: str = 'bar'):
        """
        添加图表页
        
        Args:
            title: 标题
            data: 图表数据 {'categories': [...], 'series': {'系列1': [...], ...}}
            chart_type: 图表类型 ('bar', 'line', 'pie')
        """
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb_tuple_to_color(self.style['background'])
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._rgb_tuple_to_color(self.style['primary_color'])
        
        # 图表数据
        chart_data = CategoryChartData()
        chart_data.categories = data.get('categories', [])
        
        for series_name, values in data.get('series', {}).items():
            chart_data.add_series(series_name, values)
        
        # 添加图表
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(5)
        
        chart_type_map = {
            'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE
        }
        
        chart = slide.shapes.add_chart(
            chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
            x, y, cx, cy, chart_data
        ).chart
        
        return slide
    
    def save(self, output_path: str):
        """
        保存 PPT 文件
        
        Args:
            output_path: 输出文件路径
        """
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_file))
        print(f"✅ PPT 已保存: {output_path}")
    
    def _rgb_tuple_to_color(self, rgb_tuple):
        """将 RGB 元组转换为 RGBColor"""
        from pptx.dml.color import RGBColor
        return RGBColor(*rgb_tuple)


def main():
    """测试函数"""
    # 生成三种风格的示例
    for style_name in ['conservative', 'visual', 'detailed']:
        gen = PPTGenerator(style=style_name)
        
        # 添加标题页
        gen.add_title_slide(
            f"季度业绩报告 - {gen.style['name']}", 
            "2024 Q4"
        )
        
        # 添加数据页
        sample_df = pd.DataFrame({
            '指标': ['收益率', '波动率', '夏普比率'],
            'Q3': [8.5, 12.3, 0.69],
            'Q4': [12.8, 15.1, 0.85]
        })
        gen.add_data_slide("关键指标对比", sample_df)
        
        # 添加图表页
        chart_data = {
            'categories': ['1月', '2月', '3月'],
            'series': {
                '收益率': [5.2, 7.8, 12.8],
                '基准': [4.1, 6.5, 9.2]
            }
        }
        gen.add_chart_slide("月度收益趋势", chart_data, chart_type='line')
        
        # 保存
        gen.save(f'data/output/demo_{style_name}.pptx')


if __name__ == '__main__':
    main()
